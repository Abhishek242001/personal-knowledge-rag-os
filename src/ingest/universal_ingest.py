"""
Universal Ingestion Pipeline
=============================
Walks my_knowledge/ directory, detects project folders,
and ingests all supported file types into ChromaDB.

Supported file types:
    - PDF              (.pdf)       → PyMuPDF
    - Word             (.docx)      → python-docx
    - Excel            (.xlsx .xls) → openpyxl / pandas
    - PowerPoint       (.pptx)      → python-pptx
    - Images           (.jpg .jpeg .png .webp .bmp) → Claude Vision API
    - Code             (.py .js .ts .java .cpp .c .go .rs .sh .yaml .json .toml .env)
    - Text / Markdown  (.txt .md)   → plain read
    - CSV              (.csv)       → pandas

Usage:
    python src/ingest/universal_ingest.py
    python src/ingest/universal_ingest.py --project "Person-Reid-Model"
    python src/ingest/universal_ingest.py --path my_knowledge/some_file.pdf
    python src/ingest/universal_ingest.py --clean --project "Person-Reid-Model"
"""

import os
import sys
import time
import hashlib
import argparse
import base64
import mimetypes
from pathlib import Path
from dotenv import load_dotenv

load_dotenv()

sys.path.append(str(Path(__file__).parent.parent))
from embed.embedder import get_embedder
from rag.chromadb_store import get_collection

# ── Constants ────────────────────────────────────────────────────────────────

KNOWLEDGE_DIR = Path("my_knowledge")
CHUNK_SIZE    = 400
CHUNK_OVERLAP = 50
MAX_FILE_SIZE = 25 * 1024 * 1024  # 25 MB hard limit

# File type groups
PDF_EXTS   = {".pdf"}
WORD_EXTS  = {".docx", ".doc"}
EXCEL_EXTS = {".xlsx", ".xls", ".csv"}
PPT_EXTS   = {".pptx", ".ppt"}
IMAGE_EXTS = {".jpg", ".jpeg", ".png", ".webp", ".bmp", ".gif"}
TEXT_EXTS  = {".txt", ".md", ".rst", ".log"}
CODE_EXTS  = {
    ".py", ".js", ".ts", ".jsx", ".tsx", ".java", ".cpp", ".c",
    ".h", ".go", ".rs", ".sh", ".bash", ".yaml", ".yml",
    ".json", ".toml", ".env", ".ini", ".cfg", ".sql", ".r",
    ".html", ".css", ".xml", ".tf", ".dockerfile",
}

ALL_SUPPORTED = (
    PDF_EXTS | WORD_EXTS | EXCEL_EXTS | PPT_EXTS |
    IMAGE_EXTS | TEXT_EXTS | CODE_EXTS
)

# Folders/files to always skip
SKIP_DIRS  = {".git", "__pycache__", "node_modules", ".obsidian",
              "output", "weight", "chromadb", ".env", "logs"}
SKIP_FILES = {".DS_Store", "Thumbs.db", ".gitignore"}


# ── Chunking ─────────────────────────────────────────────────────────────────

def chunk_text(text: str, chunk_size: int = CHUNK_SIZE,
               overlap: int = CHUNK_OVERLAP) -> list[str]:
    words = text.split()
    chunks, start = [], 0
    while start < len(words):
        end   = min(start + chunk_size, len(words))
        chunk = " ".join(words[start:end])
        if chunk.strip():
            chunks.append(chunk)
        start += chunk_size - overlap
    return chunks


# ── Extractors ────────────────────────────────────────────────────────────────

def extract_pdf(path: Path) -> str:
    try:
        import fitz
        doc   = fitz.open(str(path))
        pages = [page.get_text().strip() for page in doc]
        doc.close()
        return "\n\n".join(p for p in pages if p)
    except ImportError:
        print("   ⚠️  pymupdf not installed: pip install pymupdf")
        return ""
    except Exception as e:
        print(f"   ⚠️  PDF error: {e}")
        return ""


def extract_docx(path: Path) -> str:
    try:
        from docx import Document
        doc   = Document(str(path))
        lines = []
        for para in doc.paragraphs:
            if para.text.strip():
                lines.append(para.text.strip())
        # Also extract tables
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(
                    cell.text.strip() for cell in row.cells if cell.text.strip()
                )
                if row_text:
                    lines.append(row_text)
        return "\n".join(lines)
    except ImportError:
        print("   ⚠️  python-docx not installed: pip install python-docx")
        return ""
    except Exception as e:
        print(f"   ⚠️  DOCX error: {e}")
        return ""


def extract_excel(path: Path) -> str:
    try:
        import pandas as pd
        lines = []
        if path.suffix.lower() == ".csv":
            df = pd.read_csv(path, nrows=500)
            lines.append(f"CSV file: {path.name}")
            lines.append(df.to_string(index=False))
        else:
            xf = pd.ExcelFile(str(path))
            for sheet in xf.sheet_names:
                df = xf.parse(sheet, nrows=500)
                lines.append(f"Sheet: {sheet}")
                lines.append(df.to_string(index=False))
        return "\n\n".join(lines)
    except ImportError:
        print("   ⚠️  pandas/openpyxl not installed: pip install pandas openpyxl")
        return ""
    except Exception as e:
        print(f"   ⚠️  Excel/CSV error: {e}")
        return ""


def extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
        prs   = Presentation(str(path))
        lines = []
        for i, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
            if slide_texts:
                lines.append(f"Slide {i}: " + " | ".join(slide_texts))
        return "\n".join(lines)
    except ImportError:
        print("   ⚠️  python-pptx not installed: pip install python-pptx")
        return ""
    except Exception as e:
        print(f"   ⚠️  PPTX error: {e}")
        return ""


def extract_image(path: Path) -> str:
    """Use Claude Vision to describe the image content."""
    try:
        import anthropic
        api_key = os.getenv("ANTHROPIC_API_KEY", "")
        if not api_key or "your-key" in api_key:
            print("   ⚠️  No ANTHROPIC_API_KEY — skipping image")
            return ""

        with open(path, "rb") as f:
            image_data = base64.standard_b64encode(f.read()).decode("utf-8")

        mime_type, _ = mimetypes.guess_type(str(path))
        if not mime_type:
            mime_type = "image/jpeg"

        client   = anthropic.Anthropic(api_key=api_key)
        response = client.messages.create(
            model="claude-haiku-4-5",
            max_tokens=500,
            messages=[{
                "role": "user",
                "content": [
                    {
                        "type": "image",
                        "source": {
                            "type": "base64",
                            "media_type": mime_type,
                            "data": image_data,
                        },
                    },
                    {
                        "type": "text",
                        "text": (
                            "Describe this image in detail. Extract any visible text, "
                            "data, charts, diagrams, or important visual information. "
                            "Be specific and thorough — this description will be used "
                            "for search and retrieval."
                        ),
                    },
                ],
            }],
        )
        return response.content[0].text

    except ImportError:
        print("   ⚠️  anthropic not installed: pip install anthropic")
        return ""
    except Exception as e:
        print(f"   ⚠️  Image extraction error: {e}")
        return ""


def extract_text(path: Path) -> str:
    """Plain text / markdown / code files."""
    try:
        return path.read_text(encoding="utf-8", errors="ignore")
    except Exception as e:
        print(f"   ⚠️  Text read error: {e}")
        return ""


def extract_content(path: Path) -> tuple[str, str]:
    """
    Dispatch to the right extractor based on file extension.
    Returns (content_text, file_type_label).
    """
    ext = path.suffix.lower()

    if ext in PDF_EXTS:
        return extract_pdf(path),   "pdf"
    elif ext in WORD_EXTS:
        return extract_docx(path),  "word"
    elif ext in EXCEL_EXTS:
        return extract_excel(path), "excel"
    elif ext in PPT_EXTS:
        return extract_pptx(path),  "powerpoint"
    elif ext in IMAGE_EXTS:
        return extract_image(path), "image"
    elif ext in CODE_EXTS:
        return extract_text(path),  "code"
    elif ext in TEXT_EXTS:
        return extract_text(path),  "text"
    else:
        return "",                  "unknown"


# ── Project detection ─────────────────────────────────────────────────────────

def detect_project_name(file_path: Path, knowledge_dir: Path) -> str:
    """
    Determine project name from folder structure.

    my_knowledge/Person-Reid-Model/modules/foo.py → "Person-Reid-Model"
    my_knowledge/report.pdf                        → "root"
    """
    try:
        rel = file_path.relative_to(knowledge_dir)
        parts = rel.parts
        if len(parts) > 1:
            return parts[0]          # first subfolder = project name
        return "root"
    except ValueError:
        return "root"


# ── Core ingestor ─────────────────────────────────────────────────────────────

def ingest_file(path: Path, knowledge_dir: Path,
                collection, embedder, delay: float = 0.3) -> int:
    """
    Ingest a single file. Returns number of chunks stored.
    """
    # Size guard
    size = path.stat().st_size
    if size > MAX_FILE_SIZE:
        print(f"   ⏭️  Skipping (>{MAX_FILE_SIZE // (1024*1024)}MB): {path.name}")
        return 0
    if size == 0:
        return 0

    ext = path.suffix.lower()
    if ext not in ALL_SUPPORTED:
        return 0

    content, file_type = extract_content(path)
    if not content or len(content.strip()) < 30:
        print(f"   ⚠️  No usable content: {path.name}")
        return 0

    project_name = detect_project_name(path, knowledge_dir)
    rel_path     = str(path.relative_to(knowledge_dir)) if knowledge_dir in path.parents else str(path)

    # Prepend file name as context before chunking
    full_text = f"File: {path.name}\nProject: {project_name}\n\n{content}"
    chunks    = chunk_text(full_text)

    stored = 0
    for i, chunk in enumerate(chunks):
        if len(chunk.strip()) < 30:
            continue

        chunk_id = hashlib.md5(
            f"universal:{rel_path}:c{i}:{chunk[:50]}".encode()
        ).hexdigest()

        # Retry logic
        embedding = None
        for attempt in range(5):
            try:
                embedding = embedder.embed(chunk)
                break
            except Exception as e:
                err = str(e).lower()
                if "rate" in err:
                    wait = 20 * (attempt + 1)
                    print(f"   ⏳ Rate limited — waiting {wait}s...")
                    time.sleep(wait)
                else:
                    print(f"   ❌ Embed error: {e}")
                    break

        if embedding is None:
            continue

        collection.add(
            ids=[chunk_id],
            embeddings=[embedding],
            documents=[chunk],
            metadatas=[{
                "source":       "universal",
                "file_type":    file_type,
                "project_name": project_name,
                "filename":     path.name,
                "file_path":    str(path.resolve()),   # absolute path for file sending
                "rel_path":     rel_path,
                "extension":    ext,
                "chunk_index":  i,
                "file_size_kb": round(size / 1024, 1),
            }]
        )
        stored += 1
        time.sleep(delay)

    return stored


def ingest_directory(knowledge_dir: Path, project_filter: str = None,
                     embedder=None, collection=None) -> dict:
    """
    Walk knowledge_dir and ingest all supported files.
    Optionally filter to a specific project subfolder.
    Returns stats dict.
    """
    if not knowledge_dir.exists():
        print(f"❌ Directory not found: {knowledge_dir}")
        return {}

    if embedder is None:
        embedder = get_embedder()
    if collection is None:
        collection = get_collection()

    # Determine scan root
    if project_filter:
        scan_root = knowledge_dir / project_filter
        if not scan_root.exists():
            print(f"❌ Project folder not found: {scan_root}")
            return {}
        print(f"📁 Scanning project: {project_filter}")
    else:
        scan_root = knowledge_dir
        print(f"📁 Scanning: {knowledge_dir}")

    # Collect files — skip internal tool dirs but allow project subfolders
    # Key rule: skip if .git or __pycache__ etc appear as a path component
    # BUT only skip the contents of those dirs, not sibling files
    GIT_SKIP = {".git", "__pycache__", "node_modules", ".obsidian",
                "output", "weight", "chromadb", "logs"}

    all_files = []
    for f in scan_root.rglob("*"):
        if not f.is_file():
            continue

        # Get path relative to scan root and check each component
        try:
            rel = f.relative_to(scan_root)
        except ValueError:
            continue

        # Skip if any part of the relative path is a skip dir
        if any(part in GIT_SKIP for part in rel.parts):
            continue

        if f.name in SKIP_FILES:
            continue
        if f.suffix.lower() in ALL_SUPPORTED:
            all_files.append(f)

    if not all_files:
        print("⚠️  No supported files found.")
        return {}

    print(f"🔍 Found {len(all_files)} supported file(s)")
    print()

    stats = {}   # file_type → chunk count
    total = 0

    for f in sorted(all_files):
        project = detect_project_name(f, knowledge_dir)
        print(f"  [{f.suffix.upper()}] {project}/{f.name}")
        n = ingest_file(f, knowledge_dir, collection, embedder)
        if n:
            ext_label = f.suffix.lower().lstrip(".")
            stats[ext_label] = stats.get(ext_label, 0) + n
            total += n
            print(f"     ✅ {n} chunks stored")
        print()

    return {"total": total, "by_type": stats}


# ── CLI ───────────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(
        description="Universal file ingestion for Personal Knowledge OS"
    )
    parser.add_argument("--project", type=str,
                        help="Ingest only this project subfolder")
    parser.add_argument("--path",    type=str,
                        help="Ingest a single specific file")
    parser.add_argument("--dir",     type=str, default=str(KNOWLEDGE_DIR),
                        help=f"Knowledge directory (default: {KNOWLEDGE_DIR})")
    parser.add_argument("--clean",   action="store_true",
                        help="Delete existing chunks for the project before re-ingesting")
    args = parser.parse_args()

    knowledge_dir = Path(args.dir)

    print("🧠 Personal Knowledge OS — Universal Ingestion")
    print("=" * 50)

    embedder   = get_embedder()
    collection = get_collection()

    # Clean mode
    if args.clean and args.project:
        print(f"🗑️  Cleaning existing chunks for project: {args.project}")
        try:
            collection.delete(where={"project_name": args.project})
            print("   Done.")
        except Exception as e:
            print(f"   Warning: {e}")
        print()

    # Single file mode
    if args.path:
        f = Path(args.path)
        if not f.exists():
            print(f"❌ File not found: {f}")
            sys.exit(1)
        print(f"📄 Ingesting single file: {f.name}")
        n = ingest_file(f, knowledge_dir, collection, embedder)
        print(f"\n✅ Done. {n} chunks stored.")
        return

    # Directory mode
    result = ingest_directory(
        knowledge_dir,
        project_filter=args.project,
        embedder=embedder,
        collection=collection,
    )

    if result:
        print("=" * 50)
        print(f"✅ Ingestion complete — {result['total']} total chunks")
        print("By file type:")
        for ft, count in sorted(result["by_type"].items()):
            print(f"  .{ft}: {count} chunks")


if __name__ == "__main__":
    main()